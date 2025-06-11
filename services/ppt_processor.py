import os
import logging
import time
import json
import re
from typing import Dict, List, Tuple, Any
from .translator import TranslationService
import pandas as pd
from datetime import datetime
from utils.term_extractor import TermExtractor
from .translation_detector import TranslationDetector
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

# 导入PPT处理库
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    logging.error("python-pptx模块未安装，PPT处理功能将不可用")
    raise ImportError("请安装python-pptx模块以支持PPT文件处理: pip install python-pptx")

logger = logging.getLogger(__name__)

class PPTProcessor:
    """PPT文件处理器，用于从PPT文件中提取文本并进行翻译，输出仅包含翻译后的内容"""

    def __init__(self, translator: TranslationService):
        self.translator = translator
        self.use_terminology = True  # 默认使用术语库
        self.preprocess_terms = True  # 是否预处理术语（默认启用）
        self.term_extractor = TermExtractor()  # 术语提取器
        self.translation_detector = TranslationDetector()  # 翻译检测器
        self.skip_translated_content = True  # 是否跳过已翻译内容（默认启用）
        self.export_pdf = False  # 是否导出PDF
        self.output_format = "bilingual"  # 默认双语对照
        self.source_lang = "zh"  # 默认源语言为中文
        self.target_lang = "en"  # 默认目标语言为英文
        self.is_cn_to_foreign = True  # 默认翻译方向为中文→外语
        # 数学公式正则表达式模式
        self.latex_patterns = [
            r'\$\$(.*?)\$\$',  # 行间公式 $$...$$
            r'\$(.*?)\$',      # 行内公式 $...$
            r'\\begin\{equation\}(.*?)\\end\{equation\}',  # equation环境
            r'\\begin\{align\}(.*?)\\end\{align\}',        # align环境
            r'\\begin\{eqnarray\}(.*?)\\end\{eqnarray\}'   # eqnarray环境
        ]
        # PPT特定的翻译提示词
        self.ppt_prompt = "这是一个PPT演示文稿的内容，请提供简洁明了的翻译，保持专业术语的准确性，并确保翻译后的文本长度适合在幻灯片上显示。"

    def process_document(self, file_path: str, target_language: str, terminology: Dict, source_lang: str = "zh", target_lang: str = "en") -> str:
        """
        处理PPT文档翻译

        Args:
            file_path: 文件路径
            target_language: 目标语言名称（用于术语表查找）
            terminology: 术语词典
            source_lang: 源语言代码，默认为中文(zh)
            target_lang: 目标语言代码，默认为英文(en)

        Returns:
            str: 输出文件路径
        """
        # 设置翻译方向
        self.source_lang = source_lang
        self.target_lang = target_lang

        # 根据源语言和目标语言确定翻译方向
        self.is_cn_to_foreign = source_lang == "zh" or (source_lang == "auto" and target_lang != "zh")
        logger.info(f"翻译方向: {'中文→外语' if self.is_cn_to_foreign else '外语→中文'}")

        # 检查文件是否存在
        if not os.path.exists(file_path):
            logger.error(f"文件不存在: {file_path}")
            raise Exception("文件不存在")

        # 在程序根目录下创建输出目录
        output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "输出")
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # 获取文件名和扩展名
        file_name, ext = os.path.splitext(os.path.basename(file_path))
        time_stamp = datetime.now().strftime("%Y%m%d%H%M%S")
        output_path = os.path.join(output_dir, f"{file_name}_翻译版_{time_stamp}{ext}")

        # 获取目标语言的术语表
        target_terminology = terminology.get(target_language, {})
        logger.info(f"使用{target_language}术语表，包含 {len(target_terminology)} 个术语")

        # 用于存储翻译结果的列表
        translation_results = []

        # 复制并处理PPT文档
        try:
            # 检查文件是否可写入
            try:
                with open(file_path, 'rb') as f:
                    # 可以打开源文件
                    pass
            except PermissionError:
                logger.error(f"无法读取源文件，文件可能被占用: {file_path}")
                raise Exception(f"无法读取源文件，文件可能被占用。请确保文件未被其他程序（如PowerPoint）打开。")

            # 检查目标目录是否可写入
            try:
                test_file = os.path.join(output_dir, "__test_write__.tmp")
                with open(test_file, 'w') as f:
                    f.write("test")
                os.remove(test_file)
            except (PermissionError, IOError) as e:
                logger.error(f"无法写入输出目录: {output_dir}, 错误: {str(e)}")
                raise Exception(f"无法写入输出目录，请检查权限或以管理员身份运行程序。")

            # 复制PPT文档
            ppt = self._copy_presentation(file_path, output_path)
            logger.info(f"成功复制原PPT文档到: {output_path}")

            # 如果启用了术语预处理，先收集所有使用的术语
            used_terminology = {}
            if self.preprocess_terms:
                logger.info("=== PPT术语预处理已启用 ===")
                logger.info(f"翻译方向: {self.source_lang} -> {self.target_lang}")
                logger.info(f"术语库大小: {len(target_terminology)} 个术语")

                # 对于外语→中文翻译模式，预先将术语库键值对调并缓存
                if self.source_lang != "zh":
                    logger.info("外语→中文翻译模式，预先对调术语库键值并缓存...")
                    self.reversed_terminology = self._create_reversed_terminology(target_terminology)
                    logger.info(f"对调后的术语库大小: {len(self.reversed_terminology)} 个术语")
                else:
                    self.reversed_terminology = None

                # 收集所有幻灯片中的术语
                used_terminology = self._collect_terminology(ppt, target_terminology)
                logger.info(f"从PPT中提取了 {len(used_terminology)} 个术语")

                # 如果有使用的术语，导出到Excel文件
                if used_terminology:
                    self._export_used_terminology(used_terminology)

            # 处理PPT文档
            self._process_slides(ppt, target_terminology, translation_results, used_terminology)

            # 保存PPT文档
            try:
                ppt.save(output_path)
                logger.info(f"文件已保存到: {output_path}")
            except PermissionError:
                new_output_path = os.path.join(output_dir, f"{file_name}_翻译版_retry_{time_stamp}{ext}")
                logger.warning(f"保存文件失败，尝试使用新文件名: {new_output_path}")
                ppt.save(new_output_path)
                logger.info(f"文件已保存到: {new_output_path}")
                output_path = new_output_path

            # 保存JSON结果
            json_output = os.path.join(output_dir, f"{file_name}_翻译结果_{time_stamp}.json")
            with open(json_output, 'w', encoding='utf-8') as f:
                json.dump({"status": "success", "file": output_path}, f, ensure_ascii=False, indent=2)

            # 翻译完成后，导出Excel文件
            if translation_results:
                self.export_to_excel(translation_results, file_path, target_language)

            return output_path

        except Exception as e:
            logger.error(f"PPT处理过程出错: {str(e)}")
            raise

    def _copy_presentation(self, source_path: str, target_path: str) -> Presentation:
        """复制PPT文档"""
        ppt = Presentation(source_path)
        ppt.save(target_path)
        return Presentation(target_path)

    def _create_reversed_terminology(self, terminology: Dict[str, str]) -> Dict[str, str]:
        """
        创建反向术语库（外语术语: 中文术语）

        Args:
            terminology: 原始术语库 {中文术语: 外语术语}

        Returns:
            Dict[str, str]: 反向术语库 {外语术语: 中文术语}
        """
        reversed_terminology = {}

        for cn_term, foreign_term in terminology.items():
            # 确保外语术语不为空
            if foreign_term and foreign_term.strip():
                # 如果有多个中文术语对应同一个外语术语，选择最长的中文术语
                if foreign_term in reversed_terminology:
                    existing_cn_term = reversed_terminology[foreign_term]
                    if len(cn_term) > len(existing_cn_term):
                        reversed_terminology[foreign_term] = cn_term
                        logger.debug(f"术语冲突，选择更长的中文术语: {foreign_term} -> {cn_term} (替换 {existing_cn_term})")
                else:
                    reversed_terminology[foreign_term] = cn_term

        logger.info(f"成功创建反向术语库，原始术语库: {len(terminology)} 个，反向术语库: {len(reversed_terminology)} 个")
        return reversed_terminology

    def _collect_terminology(self, ppt: Presentation, terminology: Dict) -> Dict:
        """收集PPT中使用的术语"""
        used_terminology = {}

        # 遍历所有幻灯片
        for slide in ppt.slides:
            # 遍历幻灯片中的所有形状
            for shape in slide.shapes:
                # 处理文本框
                if hasattr(shape, "text") and shape.text.strip():
                    # 根据翻译方向选择不同的术语提取方法
                    if self.source_lang == "zh":
                        # 中文 → 外语
                        shape_terms = self.term_extractor.extract_terms(shape.text, terminology)
                    else:
                        # 外语 → 中文，使用缓存的反向术语库
                        if hasattr(self, 'reversed_terminology') and self.reversed_terminology:
                            # 使用缓存的反向术语库进行高效匹配
                            shape_terms = self.term_extractor.extract_foreign_terms_from_reversed_dict(shape.text, self.reversed_terminology)
                        else:
                            # 回退到原始方法
                            shape_terms = self.term_extractor.extract_foreign_terms_by_chinese_values(shape.text, terminology)
                    # 更新使用的术语词典
                    used_terminology.update(shape_terms)

                # 处理表格
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                # 根据翻译方向选择不同的术语提取方法
                                if self.source_lang == "zh":
                                    # 中文 → 外语
                                    cell_terms = self.term_extractor.extract_terms(cell.text, terminology)
                                else:
                                    # 外语 → 中文，使用缓存的反向术语库
                                    if hasattr(self, 'reversed_terminology') and self.reversed_terminology:
                                        # 使用缓存的反向术语库进行高效匹配
                                        cell_terms = self.term_extractor.extract_foreign_terms_from_reversed_dict(cell.text, self.reversed_terminology)
                                    else:
                                        # 回退到原始方法
                                        cell_terms = self.term_extractor.extract_foreign_terms_by_chinese_values(cell.text, terminology)
                                # 更新使用的术语词典
                                used_terminology.update(cell_terms)

        return used_terminology

    def _process_slides(self, ppt: Presentation, terminology: Dict, translation_results: List, used_terminology: Dict = None) -> None:
        """处理PPT中的所有幻灯片"""
        # 遍历所有幻灯片
        for i, slide in enumerate(ppt.slides):
            logger.info(f"正在处理第 {i+1} 张幻灯片，共 {len(ppt.slides)} 张")

            # 处理幻灯片中的所有形状
            for shape in slide.shapes:
                # 处理文本框
                if hasattr(shape, "text") and shape.text.strip():
                    self._process_text_shape(shape, terminology, translation_results, used_terminology)

                # 处理表格
                if shape.has_table:
                    self._process_table(shape.table, terminology, translation_results, used_terminology)

    def _process_text_shape(self, shape: Any, terminology: Dict, translation_results: List, used_terminology: Dict = None) -> None:
        """处理文本形状"""
        original_text = shape.text.strip()
        if not original_text:
            return

        # 检查是否应该跳过翻译（已翻译内容检测）
        if self.skip_translated_content:
            should_skip, reason = self.translation_detector.should_skip_translation(
                original_text, self.source_lang, self.target_lang
            )
            if should_skip:
                logger.info(f"跳过PPT文本形状翻译: {reason} - {original_text[:50]}...")
                return

        # 检查文本中是否包含数学公式
        text, formulas = self._extract_latex_formulas(original_text)

        # 翻译文本（不包含公式部分）
        try:
            # 添加PPT特定提示词以优化翻译质量
            translation_prompt = self.ppt_prompt

            if self.preprocess_terms and used_terminology:
                # 使用术语预处理方式翻译
                try:
                    # 根据翻译方向选择不同的术语处理方法
                    if self.source_lang == "zh":
                        # 中文 → 外语
                        processed_text = self.term_extractor.replace_terms_with_placeholders(text, used_terminology)
                        translated_with_placeholders = self.translator.translate_text(
                            processed_text,
                            None,
                            self.source_lang,
                            self.target_lang,
                            prompt=translation_prompt
                        )
                        translation = self.term_extractor.restore_placeholders_with_foreign_terms(translated_with_placeholders)
                    else:
                        # 外语 → 中文
                        # used_terminology已经是 {外语术语: 中文术语} 格式，直接使用
                        reverse_terminology = used_terminology

                        logger.info(f"使用术语预处理方式翻译，找到 {len(reverse_terminology)} 个匹配术语")
                        processed_text = self.term_extractor.replace_foreign_terms_with_placeholders(text, reverse_terminology)
                        translated_with_placeholders = self.translator.translate_text(
                            processed_text,
                            None,
                            self.source_lang,
                            self.target_lang,
                            prompt=translation_prompt
                        )
                        translation = self.term_extractor.restore_placeholders_with_chinese_terms(translated_with_placeholders)
                except Exception as e:
                    logger.error(f"术语预处理翻译失败: {str(e)}")
                    # 如果术语预处理失败，使用常规翻译
                    translation = self.translator.translate_text(
                        text,
                        None,
                        self.source_lang,
                        self.target_lang,
                        prompt=translation_prompt
                    )
            else:
                # 使用常规方式翻译
                translation = self.translator.translate_text(
                    text,
                    terminology,
                    self.source_lang,
                    self.target_lang,
                    prompt=translation_prompt
                )

            # 将公式重新插入到翻译后的文本中
            if formulas:
                translation = self._restore_latex_formulas(translation, formulas)

            # 记录翻译结果
            translation_results.append({
                "原文": original_text,
                "译文": translation,
                "位置": f"幻灯片 {shape.part.slide.slide_id}，形状 {shape.shape_id}"
            })

            # 根据输出格式设置文本内容
            if self.output_format == "bilingual":
                # 双语对照模式
                shape.text = f"{original_text}\n\n{translation}"
            else:
                # 仅翻译模式
                shape.text = translation

            # 尝试自动调整文本框大小以适应内容
            self._adjust_text_shape_size(shape, original_text, translation)

        except Exception as e:
            logger.error(f"翻译文本形状失败: {str(e)}")
            # 在出错时，保留原文
            shape.text = original_text

    def _adjust_text_shape_size(self, shape: Any, original_text: str, translation: str) -> None:
        """根据翻译后文本长度自动调整文本框大小"""
        try:
            # 只有当形状是文本框时才调整大小
            if hasattr(shape, "text_frame"):
                # 计算原文和译文的长度比例
                original_lines = original_text.count('\n') + 1
                translation_lines = translation.count('\n') + 1

                # 如果是双语对照模式，考虑两者的总行数
                if self.output_format == "bilingual":
                    total_lines = original_lines + translation_lines + 1  # +1 是为了中间的空行
                else:
                    total_lines = translation_lines

                # 获取当前文本框高度
                current_height = shape.height

                # 如果翻译后的文本行数明显多于原文，增加文本框高度
                if total_lines > original_lines * 1.2:  # 如果总行数超过原文的1.2倍
                    # 计算新高度，但不超过原高度的2倍
                    ratio = min(total_lines / original_lines, 2.0)
                    new_height = int(current_height * ratio)

                    # 调整文本框高度，保持顶部位置不变
                    shape.height = new_height

                    logger.info(f"调整文本框大小: 原高度={current_height}, 新高度={new_height}")
        except Exception as e:
            logger.warning(f"调整文本框大小失败: {str(e)}")
            # 调整失败不影响翻译功能，只记录警告

    def _process_table(self, table: Any, terminology: Dict, translation_results: List, used_terminology: Dict = None) -> None:
        """处理表格"""
        for row in table.rows:
            for cell in row.cells:
                original_text = cell.text.strip()
                if not original_text:
                    continue

                # 检查是否应该跳过翻译（已翻译内容检测）
                if self.skip_translated_content:
                    should_skip, reason = self.translation_detector.should_skip_translation(
                        original_text, self.source_lang, self.target_lang
                    )
                    if should_skip:
                        logger.info(f"跳过PPT表格单元格翻译: {reason} - {original_text[:50]}...")
                        continue

                # 检查文本中是否包含数学公式
                text, formulas = self._extract_latex_formulas(original_text)

                # 翻译文本（不包含公式部分）
                try:
                    # 添加PPT特定提示词以优化翻译质量
                    translation_prompt = self.ppt_prompt

                    if self.preprocess_terms and used_terminology:
                        # 使用术语预处理方式翻译
                        try:
                            # 根据翻译方向选择不同的术语处理方法
                            if self.source_lang == "zh":
                                # 中文 → 外语
                                logger.info(f"使用术语预处理方式翻译，找到 {len(used_terminology)} 个匹配术语")
                                processed_text = self.term_extractor.replace_terms_with_placeholders(text, used_terminology)
                                translated_with_placeholders = self.translator.translate_text(
                                    processed_text,
                                    None,
                                    self.source_lang,
                                    self.target_lang,
                                    prompt=translation_prompt
                                )
                                translation = self.term_extractor.restore_placeholders_with_foreign_terms(translated_with_placeholders)
                            else:
                                # 外语 → 中文
                                # used_terminology已经是 {外语术语: 中文术语} 格式，直接使用
                                reverse_terminology = used_terminology

                                logger.info(f"使用术语预处理方式翻译，找到 {len(reverse_terminology)} 个匹配术语")
                                processed_text = self.term_extractor.replace_foreign_terms_with_placeholders(text, reverse_terminology)
                                translated_with_placeholders = self.translator.translate_text(
                                    processed_text,
                                    None,
                                    self.source_lang,
                                    self.target_lang,
                                    prompt=translation_prompt
                                )
                                translation = self.term_extractor.restore_placeholders_with_chinese_terms(translated_with_placeholders)
                        except Exception as e:
                            logger.error(f"表格单元格术语预处理翻译失败: {str(e)}")
                            # 如果术语预处理失败，使用常规翻译
                            translation = self.translator.translate_text(
                                text,
                                None,
                                self.source_lang,
                                self.target_lang,
                                prompt=translation_prompt
                            )
                    else:
                        # 使用常规方式翻译
                        translation = self.translator.translate_text(
                            text,
                            terminology,
                            self.source_lang,
                            self.target_lang,
                            prompt=translation_prompt
                        )

                    # 将公式重新插入到翻译后的文本中
                    if formulas:
                        translation = self._restore_latex_formulas(translation, formulas)

                    # 记录翻译结果
                    translation_results.append({
                        "原文": original_text,
                        "译文": translation,
                        "位置": f"表格单元格"
                    })

                    # 根据输出格式设置文本内容
                    if self.output_format == "bilingual":
                        # 双语对照模式
                        cell.text = f"{original_text}\n\n{translation}"
                    else:
                        # 仅翻译模式
                        cell.text = translation

                except Exception as e:
                    logger.error(f"翻译表格单元格失败: {str(e)}")
                    # 在出错时，保留原文
                    cell.text = original_text

    def _extract_latex_formulas(self, text: str) -> Tuple[str, List[Tuple[str, str]]]:
        """
        从文本中提取LaTeX公式，并用占位符替换

        Args:
            text: 原始文本

        Returns:
            Tuple[str, List[Tuple[str, str]]]: 替换后的文本和公式列表（占位符，公式）
        """
        if not text:
            return text, []

        # 存储提取的公式
        formulas = []
        # 替换后的文本
        processed_text = text

        # 遍历所有公式模式
        for i, pattern in enumerate(self.latex_patterns):
            # 查找所有匹配的公式
            matches = re.finditer(pattern, processed_text)

            # 替换公式为占位符
            for j, match in enumerate(matches):
                formula = match.group(0)  # 完整公式，包括分隔符
                placeholder = f"[FORMULA_{i}_{j}]"
                processed_text = processed_text.replace(formula, placeholder, 1)
                formulas.append((placeholder, formula))

        return processed_text, formulas

    def _restore_latex_formulas(self, text: str, formulas: List[Tuple[str, str]]) -> str:
        """
        将公式占位符替换回原始公式

        Args:
            text: 包含占位符的文本
            formulas: 公式列表（占位符，公式）

        Returns:
            str: 恢复公式后的文本
        """
        result = text
        for placeholder, formula in formulas:
            result = result.replace(placeholder, formula)
        return result

    def _export_used_terminology(self, used_terminology: Dict) -> None:
        """导出使用的术语到Excel文件，包含使用统计信息"""
        try:
            # 创建输出目录
            output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "输出")
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)

            # 创建Excel文件
            time_stamp = datetime.now().strftime("%Y%m%d%H%M%S")
            excel_path = os.path.join(output_dir, f"使用的术语_{time_stamp}.xlsx")

            # 根据翻译方向创建不同的列名
            if self.source_lang == "zh":
                # 中文 → 外语
                source_column = '中文术语'
                target_column = '目标语术语'
            else:
                # 外语 → 中文
                source_column = '外语术语'
                target_column = '中文术语'

            # 获取术语使用统计信息
            usage_stats = self.term_extractor.get_terminology_usage_stats()

            # 如果有统计信息，使用更详细的导出格式
            if usage_stats:
                # 创建包含使用次数的DataFrame
                data = []
                for term, term_info in usage_stats.items():
                    data.append({
                        source_column: term_info['source'],
                        target_column: term_info['target'],
                        '使用次数': term_info['count']
                    })

                # 如果没有统计数据，使用基本的术语列表
                if not data:
                    for source, target in used_terminology.items():
                        data.append({
                            source_column: source,
                            target_column: target,
                            '使用次数': 'N/A'
                        })

                # 创建DataFrame
                df = pd.DataFrame(data)

                # 按使用次数降序排序
                if '使用次数' in df.columns and any(isinstance(x, int) for x in df['使用次数']):
                    df = df.sort_values(by='使用次数', ascending=False)
            else:
                # 创建基本的DataFrame（无使用次数）
                df = pd.DataFrame(list(used_terminology.items()), columns=[source_column, target_column])

            # 保存到Excel
            df.to_excel(excel_path, index=False)
            logger.info(f"使用的术语已导出到: {excel_path}")

        except Exception as e:
            logger.error(f"导出术语到Excel失败: {str(e)}")

    def export_to_excel(self, translation_results: List[Dict], file_path: str, target_language: str) -> None:
        """
        将翻译结果导出到Excel文件

        Args:
            translation_results: 翻译结果列表
            file_path: 原始文件路径
            target_language: 目标语言
        """
        try:
            # 创建输出目录
            output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "输出")
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)

            # 获取文件名
            file_name = os.path.splitext(os.path.basename(file_path))[0]
            time_stamp = datetime.now().strftime("%Y%m%d%H%M%S")
            excel_path = os.path.join(output_dir, f"{file_name}_翻译结果_{time_stamp}.xlsx")

            # 转换翻译结果为DataFrame
            df = pd.DataFrame(translation_results)

            # 保存到Excel
            df.to_excel(excel_path, index=False)
            logger.info(f"翻译结果已导出到: {excel_path}")

        except Exception as e:
            logger.error(f"导出翻译结果到Excel失败: {str(e)}")
